"""Generate the FY27 tagging and tracking PowerPoint deck.

The editable source is a Markdown file with sections named "## Slide N: Title".
This script converts those sections into editable PowerPoint text boxes.
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SOURCE = Path("FY27 Tagging and Tracking Training Deck.md")
OUTPUT = Path("FY27 Tagging and Tracking Training Deck.pptx")
FOOTER_TITLE = "FY27 Tagging and Tracking Training"

BLUE = RGBColor(0, 85, 150)
DARK = RGBColor(31, 41, 55)
GRAY = RGBColor(92, 112, 128)
LIGHT_BLUE = RGBColor(226, 240, 250)
PALE = RGBColor(245, 248, 251)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(0, 128, 96)
ORANGE = RGBColor(214, 120, 0)


def normalize_line(line: str) -> str:
    """Convert lightweight Markdown styling into plain editable slide text."""

    line = re.sub(r"\*\*(.*?)\*\*", r"\1", line)
    line = re.sub(r"`([^`]*)`", r"\1", line)
    return line.replace("&amp;", "&").rstrip()


def slide_blocks(text: str) -> list[str]:
    """Extract slide sections from the Markdown source."""

    blocks = []
    for part in re.split(r"\n---\n", text):
        section = part.strip()
        if section.startswith("## Slide "):
            blocks.append(section)
    return blocks


def split_notes(block: str) -> tuple[str, str]:
    """Return slide body and facilitator notes from a slide section."""

    marker = "\nSpeaker notes:"
    if marker in block:
        body, notes = block.split(marker, 1)
        return body.strip(), notes.strip()
    return block.strip(), ""


def body_lines(body: str) -> list[tuple[str, str, int]]:
    """Parse supported Markdown lines into simple typed text rows."""

    rows: list[tuple[str, str, int]] = []
    in_code = False
    code_lines: list[str] = []

    for raw in body.splitlines():
        line = raw.rstrip()

        if line.strip().startswith("```"):
            if in_code:
                rows.append(("code", "\n".join(code_lines), 0))
                code_lines = []
                in_code = False
            else:
                in_code = True
            continue

        if in_code:
            code_lines.append(line)
            continue

        if not line.strip():
            continue

        if line.strip().startswith("|") and line.strip().endswith("|"):
            cells = [normalize_line(cell.strip()) for cell in line.strip("|").split("|")]
            if all(set(cell) <= {"-"} for cell in cells):
                continue
            rows.append(("table", " | ".join(cells), 0))
            continue

        indent = len(raw) - len(raw.lstrip(" "))
        stripped = normalize_line(line.strip())

        if stripped.startswith("- "):
            rows.append(("bullet", stripped[2:], 1 if indent >= 2 else 0))
        elif re.match(r"^\d+\.\s+", stripped):
            rows.append(("number", re.sub(r"^\d+\.\s+", "", stripped), 0))
        else:
            rows.append(("text", stripped, 0))

    return rows


def add_footer(prs: Presentation, slide, num: int) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.12), Inches(12.45), Inches(0.22))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = f"{FOOTER_TITLE} | Repository-backed draft | Slide {num}"
    paragraph.font.size = Pt(7)
    paragraph.font.color.rgb = GRAY
    paragraph.alignment = PP_ALIGN.RIGHT


def add_title(prs: Presentation, slide, title: str) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.38), Inches(12.0), Inches(0.65))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.bold = True
    paragraph.font.size = Pt(28 if len(title) < 55 else 24)
    paragraph.font.color.rgb = BLUE


def add_text_content(slide, rows: list[tuple[str, str, int]], top=Inches(1.25), height=Inches(4.95)) -> None:
    box = slide.shapes.add_textbox(Inches(0.72), top, Inches(11.9), height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.03)
    frame.clear()

    visible_lines = sum(row[1].count("\n") + 1 if row[0] == "code" else 1 for row in rows)
    font_size = 15
    if visible_lines > 18:
        font_size = 11
    elif visible_lines > 14:
        font_size = 12
    elif visible_lines > 10:
        font_size = 13

    first = True
    header_lines = {
        "Every standard tracking URL must include:",
        "Conditional:",
        "Reserved:",
        "Required today:",
        "Examples:",
        "Important:",
        "Scenario:",
        "Correct format:",
        "Why it works:",
        "Use ccid and dtid for:",
        "Do not:",
        "Use UTMs whenever an activation drives traffic to a Cisco central marketing-owned property.",
        "Repository-backed rule:",
        "Training implication:",
        "FY27 owner edits:",
        "Before a URL is used, confirm:",
        "Open items for FY27:",
    }

    for kind, content, level in rows:
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        paragraph.space_after = Pt(3)
        paragraph.font.color.rgb = DARK

        if kind == "bullet":
            paragraph.text = content
            paragraph.level = level
            paragraph.font.size = Pt(font_size)
        elif kind == "number":
            paragraph.text = content
            paragraph.level = 0
            paragraph.font.size = Pt(font_size)
        elif kind == "table":
            paragraph.text = content
            paragraph.font.size = Pt(max(9, font_size - 2))
            paragraph.font.color.rgb = GRAY
        elif kind == "code":
            paragraph.text = content
            paragraph.font.name = "Courier New"
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = GREEN
        else:
            paragraph.text = content
            paragraph.font.size = Pt(font_size)
            if content.endswith(":") or content in header_lines:
                paragraph.font.bold = True
                paragraph.font.color.rgb = BLUE


def add_notes(slide, notes: str) -> None:
    if not notes:
        return

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.72),
        Inches(6.22),
        Inches(11.9),
        Inches(0.72),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALE
    shape.line.color.rgb = LIGHT_BLUE

    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = "Facilitator note: " + " ".join(
        normalize_line(line.strip()).lstrip("- ") for line in notes.splitlines() if line.strip()
    )
    paragraph.font.size = Pt(9)
    paragraph.font.color.rgb = GRAY


def add_title_slide(prs: Presentation, blank_layout, block: str, num: int) -> None:
    slide = prs.slides.add_slide(blank_layout)
    title_line = block.splitlines()[0]
    title = title_line.split(": ", 1)[1]

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.0))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.72), Inches(1.55), Inches(11.7), Inches(1.35))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.bold = True
    paragraph.font.size = Pt(38)
    paragraph.font.color.rgb = BLUE

    subtitle = frame.add_paragraph()
    subtitle.text = "Repository-backed draft for FY27 training documentation"
    subtitle.font.size = Pt(18)
    subtitle.font.color.rgb = GRAY

    body, notes = split_notes("\n".join(block.splitlines()[1:]))
    add_text_content(slide, body_lines(body), top=Inches(3.15), height=Inches(2.1))
    add_notes(slide, notes)
    add_footer(prs, slide, num)


def build_deck() -> None:
    source_text = SOURCE.read_text(encoding="utf-8")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for idx, block in enumerate(slide_blocks(source_text), start=1):
        if idx == 1:
            add_title_slide(prs, blank_layout, block, idx)
            continue

        heading, rest = block.split("\n", 1)
        title = heading.split(": ", 1)[1]
        body, notes = split_notes(rest)
        rows = body_lines(body)

        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE
        add_title(prs, slide, title)

        if "Precedence and persistence" in title:
            card = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(0.72),
                Inches(1.18),
                Inches(11.9),
                Inches(0.65),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = LIGHT_BLUE
            card.line.color.rgb = BLUE
            frame = card.text_frame
            frame.clear()
            paragraph = frame.paragraphs[0]
            paragraph.text = "Key repository-backed rule: utm_medium is the highest-priority channel classification input."
            paragraph.font.bold = True
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = BLUE
            add_text_content(slide, rows[2:], top=Inches(2.0), height=Inches(4.05))
        elif "What happens when tagging is wrong" in title:
            card = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(0.72),
                Inches(1.2),
                Inches(11.9),
                Inches(0.58),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 243, 224)
            card.line.color.rgb = ORANGE
            frame = card.text_frame
            frame.clear()
            paragraph = frame.paragraphs[0]
            paragraph.text = "Bad tags create bad data; bad data drives the wrong optimization decisions."
            paragraph.font.bold = True
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = ORANGE
            add_text_content(slide, rows[1:], top=Inches(1.95), height=Inches(4.1))
        else:
            add_text_content(slide, rows)

        add_notes(slide, notes)
        add_footer(prs, slide, idx)

    prs.save(OUTPUT)
    print(f"Generated {OUTPUT} with {len(prs.slides)} slides from {SOURCE}")


if __name__ == "__main__":
    build_deck()
