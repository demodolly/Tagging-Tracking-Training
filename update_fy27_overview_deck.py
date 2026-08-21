"""Apply editorial-guide updates to Tagging & Tracking Overview - FY27.pptx.

Preserves the Cisco template, visuals, and module deep dives while fixing:
- utm_id / ccid inconsistency on the hybrid URL slide
- Workfront Channel ID future-state wording
- Slide 34 decision matrix (scenario-based)
- Key placeholder / framing copy on slides 1, 4, 22, 23, 27, 30, 32
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SOURCE = Path("Tagging & Tracking Overview - FY27.pptx")
OUTPUT = Path("Tagging & Tracking Overview - FY27.pptx")
BACKUP = Path("Tagging & Tracking Overview - FY27.backup.pptx")

BLUE = RGBColor(0, 85, 150)
DARK = RGBColor(31, 41, 55)
GRAY = RGBColor(92, 112, 128)
GREEN = RGBColor(0, 128, 96)
ORANGE = RGBColor(214, 120, 0)
PALE = RGBColor(245, 248, 251)


def set_text(shape, text: str, *, size: int = 11, bold: bool = False, color: RGBColor = DARK) -> None:
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, *, size=11, bold=False, color=DARK, fill=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    if fill:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = BLUE
    set_text(box, text, size=size, bold=bold, color=color)
    return box


def find_shape_by_text(slide, needle: str):
    for shape in slide.shapes:
        if hasattr(shape, "text") and needle in shape.text:
            return shape
    return None


def update_slide_1(slide) -> None:
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip() == "Navigating the Hybrid Architecture to Automate Analytics":
            set_text(
                shape,
                "Which IDs to use, where they go, and why — during the FY27 hybrid transition",
                size=18,
                color=GRAY,
            )
            break


def update_slide_4(slide) -> None:
    must_text = (
        "On every external URL that drives traffic to Cisco.com, you MUST include today:\n\n"
        "• ccid — Activity ID (campaign / initiative in system of record)\n"
        "• dtid — Drive To ID (legacy channel + vehicle classification)\n"
        "• utm_id — must equal ccid today (future: Workfront Channel ID)\n"
        "• utm_medium — channel classification (highest priority in analytics)\n"
        "• utm_source — platform or vendor within the channel\n"
        "• utm_creative=%ecid! — only for Paid Direct, Paid Programmatic, Paid Social\n\n"
        "Do not drop ccid or dtid until governance announces retirement."
    )
    for shape in slide.shapes:
        if shape.name == "TextBox 15" or (hasattr(shape, "text") and shape.text.strip() == "" and shape.top > Inches(3)):
            set_text(shape, must_text, size=10, color=DARK)
            return
    add_textbox(slide, Inches(0.6), Inches(3.0), Inches(12.0), Inches(3.5), must_text, size=10)


def update_slide_22(slide) -> None:
    add_textbox(
        slide,
        Inches(0.55),
        Inches(6.55),
        Inches(12.2),
        Inches(0.55),
        "These four CTT IDs power lead records, nurture, and pipeline — not Adobe channel reports.",
        size=9,
        bold=True,
        color=BLUE,
        fill=PALE,
    )
    add_textbox(
        slide,
        Inches(8.8),
        Inches(4.8),
        Inches(3.8),
        Inches(0.9),
        "Event ID (EID)\n(When / where the interaction happened)",
        size=10,
        bold=True,
        color=DARK,
    )


def update_slide_23(slide) -> None:
    shape = find_shape_by_text(slide, "The Core Dilemma")
    if shape:
        existing = shape.text.strip()
        if "not replacing CTT overnight" not in existing:
            shape.text_frame.paragraphs[0].text = existing + (
                "\n\nThis is why we are adding UTMs — not replacing CTT overnight."
            )
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(14)


def update_slide_27(slide) -> None:
    steps = (
        "1. Marketer selects channel type and platform in Workfront.\n"
        "2. Workfront assigns CCID, DTID, and maps utm_medium / utm_source from Source and Mediums.xlsx.\n"
        "3. Generated query string is appended to the destination URL — no manual UTM editing.\n\n"
        "Use Stensul for supported email flows; Manual URL Builder for web-referral and exceptions.\n"
        "If a source or medium is missing → escalate. Do not create local variants."
    )
    add_textbox(slide, Inches(0.55), Inches(5.6), Inches(12.2), Inches(1.5), steps, size=10, color=DARK, fill=PALE)


def update_slide_28(slide) -> None:
    mapping = {
        "TextBox 44": "https://www.cisco.com/site/us/en/about/why-cisco/ai-ready-data-centers/index.html?ccid=cc010375",
        "TextBox 45": "&dtid=pdixsp001642",
        "TextBox 47": "&utm_id=cc010375",
        "TextBox 48": "&utm_medium=paid-direct",
        "TextBox 49": "&utm_source=businessinsider&utm_creative=%ecid!",
    }
    for shape in slide.shapes:
        if shape.name in mapping:
            set_text(shape, mapping[shape.name], size=10, color=GREEN if shape.name != "TextBox 44" else DARK)

    label = find_shape_by_text(slide, "We are in an active transition")
    if label:
        set_text(
            label,
            "TODAY — Required hybrid URL (FY27): utm_id must match ccid. "
            "FUTURE preview (do not use until governance confirms): utm_id=CHL000093 with ccid/dtid removed.",
            size=10,
            bold=True,
            color=BLUE,
        )

    add_textbox(
        slide,
        Inches(0.55),
        Inches(6.65),
        Inches(12.2),
        Inches(0.45),
        "✓ utm_id matches ccid  ·  all values lowercase  ·  & separates parameters",
        size=9,
        bold=True,
        color=GREEN,
    )


def update_slide_30(slide) -> None:
    add_textbox(
        slide,
        Inches(0.55),
        Inches(6.7),
        Inches(12.2),
        Inches(0.45),
        "Workfront IDs govern intake and automation; CTT IDs still populate lead context until retirement.",
        size=10,
        bold=True,
        color=BLUE,
        fill=PALE,
    )


def update_slide_31(slide) -> None:
    for shape in slide.shapes:
        if shape.name == "TextBox 19" or (
            hasattr(shape, "text") and "MUST be appended to the utm_id field" in shape.text
        ):
            set_text(
                shape,
                "Today: utm_id must equal ccid on every external URL. "
                "Future: Workfront Channel ID will replace ccid in utm_id once governance confirms go-live.",
                size=9,
                color=DARK,
            )
            break


def update_slide_32(slide) -> None:
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.startswith("Structure of an Channel ID"):
            set_text(shape, "Structure of a Channel ID", size=24, bold=True, color=BLUE)
            break


def replace_slide_34_table(slide) -> None:
    table_shape = None
    for shape in slide.shapes:
        if shape.has_table:
            table_shape = shape
            break
    if not table_shape:
        return

    left, top, width, height = table_shape.left, table_shape.top, table_shape.width, table_shape.height
    sp = table_shape._element
    sp.getparent().remove(sp)

    rows, cols = 7, 4
    new_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = new_shape.table

    headers = ["Scenario", "Put on the URL", "Put elsewhere", "Do NOT"]
    data = [
        (
            "External link to Cisco.com (email, paid, social, syndication, referral)",
            "ccid, dtid, utm_id, utm_medium, utm_source (+ utm_creative if paid)",
            "Workfront Channel ID auto-mapped when live",
            "Use Offer ID as a URL param",
        ),
        (
            "Internal Cisco.com CTA or page link",
            "Nothing",
            "—",
            "Any CTT or UTM params",
        ),
        (
            "Gated offer landing page (AEM)",
            "—",
            "Offer ID + Content ID in page HTML / AEM metadata",
            "Tag internal navigation with DTID",
        ),
        (
            "Manual upload (MUSE template)",
            "—",
            "CCID + DTID in template columns; Offer/Event as applicable",
            "UTM params (not a live URL)",
        ),
        (
            "Integrate / PathFactory / BrightTalk",
            "Per integration spec",
            "CCID, DTID, Offer ID in payload",
            "Assume UTMs replace integration metadata",
        ),
    ]
    rule = (
        "Rule of thumb: Outside traffic → hybrid URL (let Workfront generate). "
        "Lead needs seller context → CTT on the record. "
        "Gated asset → Offer ID on the page. "
        "Project ID → internal only."
    )

    all_rows = [headers] + list(data) + [("Quick rule of thumb", rule, "", "")]

    for ri, row_data in enumerate(all_rows):
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8 if ri > 0 else 9)
                paragraph.font.color.rgb = DARK
                if ri == 0 or (ri == len(all_rows) - 1 and ci == 0):
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = BLUE
                if ri == len(all_rows) - 1 and ci == 1:
                    paragraph.font.size = Pt(8)
                    paragraph.font.color.rgb = ORANGE


def build() -> None:
    if not SOURCE.exists():
        raise FileNotFoundError(f"Missing source deck: {SOURCE}")

    if not BACKUP.exists():
        BACKUP.write_bytes(SOURCE.read_bytes())

    prs = Presentation(str(SOURCE))

    updaters = {
        1: update_slide_1,
        4: update_slide_4,
        22: update_slide_22,
        23: update_slide_23,
        27: update_slide_27,
        28: update_slide_28,
        30: update_slide_30,
        31: update_slide_31,
        32: update_slide_32,
        34: replace_slide_34_table,
    }

    for num, fn in updaters.items():
        fn(prs.slides[num - 1])

    prs.save(str(OUTPUT))
    print(f"Updated {OUTPUT} ({len(prs.slides)} slides). Backup: {BACKUP}")


if __name__ == "__main__":
    build()
