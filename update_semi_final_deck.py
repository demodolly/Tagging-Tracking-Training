"""Align Tagging & Tracking Overview - FY27 - Semi Final.pptx with stated Purpose and Outcome.

Purpose: help users create tracking URLs without confusion; learn which tags/UTMs are
needed right now; understand how daily data entry affects downstream reporting.

Outcome: confidence to build, validate, and deploy flawless URLs; correct lead credit;
fractional attribution; accurate seller context.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

SOURCE = Path("Tagging & Tracking Overview - FY27 - Semi Final.pptx")
OUTPUT = Path("Tagging & Tracking Overview - FY27 - Semi Final.pptx")
BACKUP = Path("Tagging & Tracking Overview - FY27 - Semi Final.backup.pptx")

BLUE = RGBColor(0, 85, 150)
DARK = RGBColor(31, 41, 55)
GRAY = RGBColor(92, 112, 128)
GREEN = RGBColor(0, 128, 96)
ORANGE = RGBColor(214, 120, 0)
PALE = RGBColor(245, 248, 251)


def set_text(shape, text: str, *, size: int = 11, bold: bool = False, color: RGBColor = DARK) -> None:
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, *, size=11, bold=False, color=DARK, fill=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    if fill:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = BLUE
    set_text(box, text, size=size, bold=bold, color=color)
    return box


def find_shape_by_text(slide, needle: str, *, partial: bool = True):
    needle_lower = needle.lower()
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            hay = shape.text.lower()
            if (partial and needle_lower in hay) or (not partial and hay == needle_lower):
                return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                if hasattr(sub, "text") and sub.text.strip():
                    hay = sub.text.lower()
                    if (partial and needle_lower in hay) or (not partial and hay == needle_lower):
                        return sub
    return None


def set_group_text_containing(slide, needle: str, new_text: str) -> bool:
    needle_lower = needle.lower()
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.GROUP:
            continue
        for sub in shape.shapes:
            if hasattr(sub, "text") and needle_lower in sub.text.lower():
                sub.text = new_text
                return True
    return False


def table_comments(slide) -> list:
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        last_col = len(table.columns) - 1
        cells = []
        for row in table.rows:
            idx = min(last_col, len(row.cells) - 1)
            cells.append(row.cells[idx])
        return cells
    return []


def update_slide_1(slide) -> None:
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip() == "Navigating the Hybrid Architecture to Automate Analytics":
            set_text(
                shape,
                "Create, validate, and deploy tracking URLs with confidence — "
                "knowing which tags and UTMs your campaigns need right now",
                size=16,
                color=GRAY,
            )
            break


def update_slide_3(slide) -> None:
    bridge = (
        "Every field you enter in Workfront or the URL Builder flows into lead credit, "
        "fractional attribution, pipeline reporting, and the context sellers see in CRM."
    )
    for shape in slide.shapes:
        if shape.name == "TextBox 29" or (
            hasattr(shape, "text") and shape.text.strip() == "" and shape.top > Inches(4)
        ):
            set_text(shape, bridge, size=11, bold=True, color=BLUE)
            return
    add_textbox(slide, Inches(0.55), Inches(5.8), Inches(12.2), Inches(0.7), bridge, size=11, bold=True, color=BLUE, fill=PALE)


def update_slide_4(slide) -> None:
    must = (
        "On every external URL driving traffic to Cisco.com, include TODAY:\n\n"
        "• ccid (Activity ID) — campaign initiative; required for lead creation\n"
        "• dtid (Drive To ID) — channel/vehicle that drove the click\n"
        "• utm_id — must equal ccid today (future: Workfront Channel ID)\n"
        "• utm_medium — channel classification (highest priority in analytics)\n"
        "• utm_source — platform or vendor within the channel\n"
        "• utm_creative=%ecid! — only for CM360 Paid Social, Paid Direct, Programmatic\n\n"
        "Incorrect or missing values break MSP, Adobe, and seller reporting."
    )
    for shape in slide.shapes:
        if shape.name == "TextBox 15":
            set_text(shape, must, size=10)
            return
    add_textbox(slide, Inches(0.55), Inches(2.9), Inches(7.5), Inches(3.8), must, size=10)


def update_slide_6(slide) -> None:
    add_textbox(
        slide,
        Inches(0.55),
        Inches(6.55),
        Inches(12.2),
        Inches(0.55),
        "Your daily tagging choices → lead records → fractional attribution → seller context → executive reporting",
        size=10,
        bold=True,
        color=BLUE,
        fill=PALE,
    )


def update_slide_28(slide) -> None:
    mapping = {
        "TextBox 44": "https://www.cisco.com/site/us/en/about/why-cisco/ai-ready-data-centers/index.html?ccid=cc010375",
        "TextBox 45": "&dtid=pdixsp001642",
        "TextBox 47": "&utm_id=cc010375",
        "TextBox 48": "&utm_medium=paid-direct",
        "TextBox 49": "&utm_source=businessinsider&utm_creative=%ecid!",
    }
    for shape in slide.shapes:
        if shape.name in mapping:
            set_text(shape, mapping[shape.name], size=10, color=GREEN if shape.name != "TextBox 44" else DARK)

    label = find_shape_by_text(slide, "We are in an active transition")
    if label:
        set_text(
            label,
            "TODAY (FY27): utm_id must match ccid on every live URL. "
            "FUTURE preview only — do not deploy until governance confirms: utm_id=CHL000093 with ccid/dtid retired.",
            size=10,
            bold=True,
            color=BLUE,
        )

    add_textbox(
        slide,
        Inches(0.55),
        Inches(6.65),
        Inches(12.2),
        Inches(0.45),
        "Validate before launch: utm_id = ccid · all values lowercase · no Offer ID on inbound URLs",
        size=9,
        bold=True,
        color=GREEN,
    )


def update_slide_31(slide) -> None:
    shape = find_shape_by_text(slide, "MUST be appended to the utm_id field")
    if shape:
        set_text(
            shape,
            "Today: Workfront auto-generates ccid, dtid, and UTMs — utm_id on the live URL must equal ccid. "
            "Future: Workfront Channel ID will populate utm_id once governance confirms go-live.",
            size=9,
        )


def update_slide_32(slide) -> None:
    set_group_text_containing(
        slide,
        "appended as a value in UTM_ID",
        "automatically generated at channel task creation. "
        "Today it maps to ccid on your URL; in future it will populate utm_id directly.",
    )


def update_slide_34(slide) -> None:
    add_textbox(
        slide,
        Inches(0.55),
        Inches(1.35),
        Inches(12.2),
        Inches(1.1),
        "From this module you will be able to choose the right IDs, build the URL, and validate it before launch — "
        "so every lead is credited correctly and sellers receive accurate customer context.",
        size=14,
        bold=True,
        color=DARK,
    )
    add_textbox(
        slide,
        Inches(0.55),
        Inches(2.55),
        Inches(12.2),
        Inches(0.9),
        "Choose IDs (Legacy · UTM · Workfront)  →  Build URL (Workfront first; Manual URL Builder if needed)  →  "
        "Validate checklist  →  Deploy",
        size=12,
        bold=True,
        color=BLUE,
        fill=PALE,
    )


def update_slide_35(slide) -> None:
    comments = table_comments(slide)
    if not comments:
        return
    updates = {
        3: "Manual Lead Uploads: CCID is mandatory for SFDC lead creation. DTID identifies acquisition channel. "
        "Include Offer ID when the upload relates to a specific gated asset; Event ID when tied to an event/MRF.",
        5: "PathFactory Registrations: CCID, DTID, and Offer ID required in integration payload. "
        "Event ID optional — use when linking to event/MRF context.",
        6: "BrightTalk Integration: CCID, DTID, and Offer ID required. Event ID optional for webinar/event linkage.",
    }
    for idx, text in updates.items():
        comments[idx].text = text


def update_slide_36(slide) -> None:
    comments = table_comments(slide)
    if not comments:
        return
    comments[0].text = (
        "TODAY (FY27): Every external URL needs ccid, dtid, utm_id (= ccid), utm_medium, and utm_source. "
        "Add utm_creative=%ecid! for CM360 Paid Social, Paid Direct, and Programmatic. "
        "FUTURE: utm_id will adopt Workfront Channel ID."
    )
    comments[5].text = (
        "PathFactory Registrations: Same hybrid URL rules as other inbound traffic — "
        "utm_id must equal ccid today; utm_medium and utm_source required."
    )


def update_slide_37(slide) -> None:
    comments = table_comments(slide)
    if not comments:
        return
    comments[0].text = (
        "TODAY: Workfront creates Channel ID at intake and auto-builds the hybrid URL (ccid + dtid + UTMs). "
        "Do not manually paste Channel ID into utm_id — the live URL uses utm_id=ccid. "
        "FUTURE: Channel ID will populate utm_id."
    )
    comments[5].text = (
        "PathFactory: Workfront Channel ID links the experience to your project for MSP reporting. "
        "URL still uses utm_id=ccid today."
    )


def update_slide_38(slide) -> None:
    checklist = (
        "Pre-Launch Validation Checklist\n"
        "☐ Built via Workfront (preferred) or Manual URL Builder\n"
        "☐ ccid, dtid, utm_id, utm_medium, utm_source present on external URLs\n"
        "☐ utm_id exactly matches ccid\n"
        "☐ utm_medium and utm_source are approved values for the channel\n"
        "☐ utm_creative=%ecid! included only when CM360 activates paid placements\n"
        "☐ Offer ID on gated page HTML only — not on inbound URL\n"
        "☐ No tracking params on internal Cisco.com links\n"
        "☐ Missing vendor/channel escalated before launch — never invent values"
    )
    add_textbox(
        slide,
        Inches(0.55),
        Inches(1.15),
        Inches(12.2),
        Inches(2.35),
        checklist,
        size=10,
        color=DARK,
        fill=PALE,
    )
    add_textbox(
        slide,
        Inches(0.55),
        Inches(3.55),
        Inches(12.2),
        Inches(0.45),
        "Pass this checklist → deploy with confidence that downstream data and seller context stay intact.",
        size=10,
        bold=True,
        color=GREEN,
    )


def fix_example_urls(slide) -> None:
    for shape in slide.shapes:
        if hasattr(shape, "text") and "https:////" in shape.text:
            set_text(shape, shape.text.replace("https:////", "https://"), size=9, color=GREEN)


def build() -> None:
    if not SOURCE.exists():
        raise FileNotFoundError(SOURCE)

    if not BACKUP.exists():
        BACKUP.write_bytes(SOURCE.read_bytes())

    prs = Presentation(str(SOURCE))
    updaters = {
        1: update_slide_1,
        3: update_slide_3,
        4: update_slide_4,
        6: update_slide_6,
        28: update_slide_28,
        31: update_slide_31,
        32: update_slide_32,
        34: update_slide_34,
        35: update_slide_35,
        36: update_slide_36,
        37: update_slide_37,
        38: update_slide_38,
        40: fix_example_urls,
        41: fix_example_urls,
        42: fix_example_urls,
    }

    for num, fn in updaters.items():
        fn(prs.slides[num - 1])

    prs.save(str(OUTPUT))
    print(f"Updated {OUTPUT} ({len(prs.slides)} slides). Backup: {BACKUP}")


if __name__ == "__main__":
    build()
